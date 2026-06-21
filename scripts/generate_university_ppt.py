"""Generate an editable university presentation for the CUSB RAG project."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path("reports/CUSB_AI_Chatbot_University_Presentation.pptx")

W = Inches(13.333)
H = Inches(7.5)

NAVY = RGBColor(23, 43, 77)
MAROON = RGBColor(139, 30, 22)
GOLD = RGBColor(217, 164, 65)
TEAL = RGBColor(47, 111, 109)
INK = RGBColor(37, 47, 62)
MUTED = RGBColor(92, 106, 122)
LIGHT = RGBColor(246, 248, 250)
PALE_BLUE = RGBColor(232, 239, 246)
PALE_GOLD = RGBColor(251, 245, 229)
PALE_TEAL = RGBColor(231, 243, 242)
PALE_RED = RGBColor(249, 235, 233)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(47, 133, 90)
AMBER = RGBColor(191, 126, 29)
RED = RGBColor(184, 61, 52)


def add_text(slide, x, y, w, h, text, size=18, color=INK, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font="Aptos"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, fill=WHITE, line=PALE_BLUE, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def add_circle(slide, x, y, d, fill=TEAL, line=TEAL):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    return shape


def add_line(slide, x1, y1, x2, y2, color=MUTED, width=1.5):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def set_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def footer(slide, number):
    add_line(slide, 0.55, 7.14, 12.78, 7.14, PALE_BLUE, 0.7)
    add_text(slide, 0.58, 7.2, 8.3, 0.18, "CUSB AI Chatbot | University Research Presentation", 8, MUTED)
    add_text(slide, 12.3, 7.2, 0.45, 0.18, str(number), 8, MUTED, align=PP_ALIGN.RIGHT)


def title(slide, heading, sub=None, number=None):
    add_text(slide, 0.64, 0.42, 11.9, 0.58, heading, 27, NAVY, True)
    add_rect(slide, 0.65, 1.08, 1.05, 0.05, MAROON, MAROON)
    if sub:
        add_text(slide, 1.84, 1.0, 10.5, 0.25, sub, 10, MUTED)
    if number:
        footer(slide, number)


def bullets(slide, items, x, y, w, h, size=17, color=INK, bullet_color=MAROON, gap=4):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(gap)
        p.level = 0
        p.text = "\u2022 " + item
    return box


def pill(slide, x, y, w, text, fill=PALE_TEAL, color=TEAL):
    add_rect(slide, x, y, w, 0.34, fill, fill, True)
    add_text(slide, x + 0.08, y + 0.07, w - 0.16, 0.18, text, 10, color, True, align=PP_ALIGN.CENTER)


def metric_card(slide, x, y, w, value, label, accent=TEAL, note=None):
    add_rect(slide, x, y, w, 1.18 if note else 1.0, WHITE, PALE_BLUE, True)
    add_rect(slide, x, y, 0.08, 1.18 if note else 1.0, accent, accent)
    add_text(slide, x + 0.22, y + 0.15, w - 0.35, 0.38, value, 24, accent, True)
    add_text(slide, x + 0.22, y + 0.61, w - 0.35, 0.2, label, 11, INK, True)
    if note:
        add_text(slide, x + 0.22, y + 0.87, w - 0.35, 0.18, note, 9, MUTED)


def process_box(slide, x, y, w, h, heading, detail, fill=WHITE, accent=TEAL):
    add_rect(slide, x, y, w, h, fill, PALE_BLUE, True)
    add_rect(slide, x, y, w, 0.08, accent, accent)
    add_text(slide, x + 0.12, y + 0.18, w - 0.24, 0.3, heading, 14, NAVY, True)
    add_text(slide, x + 0.12, y + 0.55, w - 0.24, h - 0.65, detail, 10, MUTED)


def bar(slide, x, y, w, label, value, total=1000, color=TEAL):
    add_text(slide, x, y, 2.2, 0.2, label, 11, INK, True)
    add_rect(slide, x + 2.25, y + 0.03, w, 0.16, PALE_BLUE, PALE_BLUE, True)
    add_rect(slide, x + 2.25, y + 0.03, w * value / total, 0.16, color, color, True)
    add_text(slide, x + 2.25 + w + 0.1, y, 0.8, 0.2, str(value), 11, color, True)


def add_slide(prs, number, heading, sub=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, heading, sub, number)
    return slide


def generate():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    prs.core_properties.title = "CUSB AI Chatbot University Presentation"
    prs.core_properties.subject = "Citation-grounded bilingual hybrid RAG system"
    prs.core_properties.author = "CUSB RAG System"

    # 1 Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_rect(slide, 0, 0, 13.333, 7.5, NAVY, NAVY)
    add_rect(slide, 0, 0, 0.22, 7.5, MAROON, MAROON)
    add_circle(slide, 10.78, 0.65, 1.6, MAROON, MAROON)
    add_circle(slide, 11.17, 1.04, 0.82, GOLD, GOLD)
    add_text(slide, 0.85, 0.82, 8.9, 0.5, "CUSB AI CHATBOT", 15, GOLD, True)
    add_text(slide, 0.85, 1.48, 10.9, 1.8,
             "A Citation-Grounded Bilingual\nHybrid RAG System",
             34, WHITE, True)
    add_text(slide, 0.88, 3.58, 9.3, 0.65,
             "for University Student Support at\nCentral University of South Bihar",
             20, RGBColor(222, 230, 239))
    add_rect(slide, 0.88, 4.78, 5.5, 0.06, GOLD, GOLD)
    add_text(slide, 0.88, 5.16, 7.7, 0.45, "University Research Presentation", 17, WHITE, True)
    add_text(slide, 0.88, 5.66, 8.0, 0.34, "Domain-specific RAG | English + Hinglish | Source-backed answers", 12, RGBColor(210, 220, 232))
    add_text(slide, 0.88, 6.88, 8.0, 0.22, "Editable PowerPoint deck generated from verified project artifacts", 9, RGBColor(189, 204, 220))
    add_text(slide, 11.72, 6.9, 0.8, 0.25, "01", 14, GOLD, True, align=PP_ALIGN.RIGHT)

    # 2 Agenda
    slide = add_slide(prs, 2, "Presentation Roadmap", "From research problem to measurable prototype")
    agenda = [
        "Introduction and Motivation", "Literature Review", "Problem Statement",
        "Research Objectives", "Proposed Methodology", "Dataset Description",
        "Analysis and Results", "Conclusion and Future Work", "References",
    ]
    for i, item in enumerate(agenda):
        col, row = i % 3, i // 3
        x, y = 0.8 + col * 4.15, 1.62 + row * 1.44
        add_circle(slide, x, y, 0.48, MAROON if i < 6 else TEAL, MAROON if i < 6 else TEAL)
        add_text(slide, x, y + 0.12, 0.48, 0.2, f"{i+1:02}", 11, WHITE, True, align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.66, y + 0.1, 3.22, 0.34, item, 14, NAVY, True)

    # 3 Intro
    slide = add_slide(prs, 3, "Introduction", "Why a university-specific assistant matters")
    process_box(slide, 0.78, 1.55, 5.52, 4.8, "Fragmented Information",
                "CUSB information is distributed across university web pages, PDFs, faculty profiles, admission bulletins, fee documents, notices and syllabus files.",
                PALE_GOLD, MAROON)
    process_box(slide, 7.03, 1.55, 5.52, 4.8, "Student-Centered Assistance",
                "Students naturally ask short, mixed-language questions about admission, programmes, faculty, fees, hostel, facilities, examination and results. The proposed assistant retrieves evidence and returns concise source-backed answers.",
                PALE_TEAL, TEAL)
    add_text(slide, 1.06, 4.2, 4.9, 0.9, "Problem: information exists,\nbut students must search manually.", 21, MAROON, True)
    add_text(slide, 7.3, 4.2, 4.9, 0.9, "Goal: ask once and receive\na verifiable answer.", 21, TEAL, True)

    # 4 Literature
    slide = add_slide(prs, 4, "Literature Review", "Established foundations used by the proposed system")
    refs = [
        ("2023", "Self-RAG | Asai et al.", "Adaptive retrieval and self-reflection improve factuality.", MAROON),
        ("2023", "RAGAS | Es et al.", "Reference-free evaluation for RAG pipelines.", TEAL),
        ("2024", "CRAG | Yan et al.", "Retrieval evaluation and corrective actions improve robustness.", GOLD),
        ("2024", "RAGCHECKER | Ru et al.", "Fine-grained diagnosis of retrieval and generation modules.", NAVY),
        ("2025", "RAG Evaluation Survey | Gan et al.", "Evaluation covers factuality, safety and efficiency.", GREEN),
        ("2026", "RAGe | Guder et al.", "Domain-specific benchmarking considers quality and resources.", MAROON),
    ]
    for i, (tag, author, desc, color) in enumerate(refs):
        y = 1.38 + i * 0.76
        pill(slide, 0.82, y, 1.15, tag, color, WHITE)
        add_text(slide, 2.18, y + 0.02, 3.25, 0.22, author, 13, NAVY, True)
        add_text(slide, 5.08, y + 0.02, 6.8, 0.22, desc, 13, INK)
    add_rect(slide, 0.82, 6.1, 11.7, 0.55, PALE_RED, PALE_RED, True)
    add_text(slide, 1.02, 6.26, 11.28, 0.2,
             "Research gap: practical bilingual, citation-grounded, university-specific assistance with systematic evaluation.",
             12, MAROON, True)

    # 5 Problem
    slide = add_slide(prs, 5, "Problem Statement", "A precise research question")
    issues = [
        ("01", "Heterogeneous sources", "Pages, PDFs, notices and profile records are scattered."),
        ("02", "Natural query variation", "Students ask in English, Hinglish, short phrases and typos."),
        ("03", "Trust and verification", "LLM-only answers may be unsupported or difficult to verify."),
    ]
    for i, (num, head, desc) in enumerate(issues):
        x = 0.86 + i * 4.18
        add_rect(slide, x, 1.72, 3.65, 2.2, WHITE, PALE_BLUE, True)
        add_text(slide, x + 0.18, 1.94, 0.5, 0.3, num, 20, MAROON, True)
        add_text(slide, x + 0.18, 2.47, 3.2, 0.28, head, 15, NAVY, True)
        add_text(slide, x + 0.18, 2.91, 3.18, 0.66, desc, 12, MUTED)
    add_rect(slide, 0.86, 4.58, 12.03, 1.45, NAVY, NAVY, True)
    add_text(slide, 1.12, 4.89, 11.5, 0.72,
             "Research Question\nCan a bilingual hybrid RAG assistant deliver useful, source-backed university answers with measurable quality?",
             19, WHITE, True, align=PP_ALIGN.CENTER)

    # 6 Objectives
    slide = add_slide(prs, 6, "Research Objectives", "Design goals and measurable outcomes")
    objectives = [
        "Build a CUSB-specific knowledge base", "Support English and Hinglish queries",
        "Combine semantic and keyword retrieval", "Return source citation metadata",
        "Reduce repeated API usage", "Handle transient API failures",
        "Evaluate answer and retrieval quality", "Deliver frontend, admin and analytics",
    ]
    for i, item in enumerate(objectives):
        col, row = i % 2, i // 2
        x, y = 0.85 + col * 6.12, 1.48 + row * 1.24
        add_circle(slide, x, y, 0.42, TEAL if col else MAROON, TEAL if col else MAROON)
        add_text(slide, x, y + 0.11, 0.42, 0.18, str(i + 1), 10, WHITE, True, align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.62, y + 0.08, 5.24, 0.25, item, 14, NAVY, True)
        add_line(slide, x + 0.62, y + 0.47, x + 5.66, y + 0.47, PALE_BLUE, 0.8)

    # 7 Architecture
    slide = add_slide(prs, 7, "Proposed System Architecture", "End-to-end implemented workflow")
    stages = [
        ("Sources", "CUSB pages\nPDFs\nProfiles", MAROON),
        ("Ingestion", "Extract\nClean\nChunk", GOLD),
        ("Indexing", "Embeddings\nBM25 index", TEAL),
        ("Retrieval", "Dense + BM25\nRRF fusion", NAVY),
        ("Answering", "Direct answer\nor grounded LLM", GREEN),
        ("Delivery", "FastAPI\nNext.js\nAdmin", MAROON),
    ]
    for i, (head, detail, color) in enumerate(stages):
        x = 0.55 + i * 2.13
        process_box(slide, x, 2.2, 1.66, 2.35, head, detail, WHITE, color)
        if i < len(stages) - 1:
            add_line(slide, x + 1.67, 3.38, x + 2.07, 3.38, GOLD, 2.2)
            add_text(slide, x + 1.77, 3.17, 0.2, 0.2, ">", 14, GOLD, True, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.55, 5.23, 12.25, 0.66, PALE_TEAL, PALE_TEAL, True)
    add_text(slide, 0.82, 5.42, 11.72, 0.22,
             "Evidence flows from public sources to retrievable chunks, then to language-aware answers with citations.",
             13, TEAL, True, align=PP_ALIGN.CENTER)

    # 8 Ingestion
    slide = add_slide(prs, 8, "Data Ingestion Methodology", "From public content to searchable evidence")
    steps = [
        ("1", "Collect", "Web pages, PDFs,\nprofiles and notices"),
        ("2", "Extract", "HTML and PDF\ntext extraction"),
        ("3", "Clean", "Remove noise and\npreserve metadata"),
        ("4", "Chunk", "Create smaller\nretrievable units"),
        ("5", "Embed", "Generate multilingual\nvectors"),
        ("6", "Index", "Build searchable\nartifacts"),
    ]
    for i, (n, head, desc) in enumerate(steps):
        x = 0.67 + i * 2.08
        add_circle(slide, x + 0.57, 1.78, 0.48, MAROON if i < 3 else TEAL, MAROON if i < 3 else TEAL)
        add_text(slide, x + 0.57, 1.9, 0.48, 0.17, n, 11, WHITE, True, align=PP_ALIGN.CENTER)
        if i < 5:
            add_line(slide, x + 1.1, 2.02, x + 2.0, 2.02, GOLD, 2)
        add_rect(slide, x, 2.57, 1.68, 2.13, WHITE, PALE_BLUE, True)
        add_text(slide, x + 0.13, 2.88, 1.42, 0.25, head, 15, NAVY, True, align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.12, 3.47, 1.44, 0.68, desc, 11, MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, 0.83, 5.39, 11.45, 0.45,
             "Each chunk retains source title, URL, category, department, file and citation metadata.",
             15, TEAL, True, align=PP_ALIGN.CENTER)

    # 9 Dataset
    slide = add_slide(prs, 9, "Dataset Description", "Indexed CUSB knowledge corpus")
    metric_card(slide, 0.76, 1.45, 2.36, "30,384", "indexed chunks", TEAL)
    metric_card(slide, 3.36, 1.45, 2.36, "1,302.8", "avg. characters / chunk", MAROON)
    metric_card(slide, 5.96, 1.45, 2.36, "39.58 M", "total text characters", NAVY)
    metric_card(slide, 8.56, 1.45, 1.88, "30,384", "vectors", GREEN)
    metric_card(slide, 10.68, 1.45, 1.88, "384", "dimensions", GOLD)
    add_rect(slide, 0.76, 3.13, 11.8, 2.73, WHITE, PALE_BLUE, True)
    add_text(slide, 1.02, 3.4, 3.1, 0.28, "Source Families", 16, NAVY, True)
    families = [
        "University info", "Administration", "Admissions and fees", "Faculty",
        "Facilities and infrastructure", "Syllabus", "Academic PDFs", "Research corpus",
    ]
    for i, family in enumerate(families):
        col, row = i % 4, i // 4
        pill(slide, 1.02 + col * 2.85, 4.02 + row * 0.66, 2.42, family, PALE_TEAL if row else PALE_GOLD, TEAL if row else MAROON)
    add_text(slide, 1.02, 5.51, 10.9, 0.18,
             "Embedding model: paraphrase-multilingual-MiniLM-L12-v2 | QA records excluded from reported indexed chunk metadata",
             10, MUTED)

    # 10 Retrieval
    slide = add_slide(prs, 10, "Hybrid Retrieval Method", "Dense semantics plus sparse keyword matching")
    process_box(slide, 0.88, 1.55, 3.1, 2.14, "Dense Retrieval",
                "Multilingual semantic vectors searched through FAISS. Local Qdrant storage is configurable.", PALE_TEAL, TEAL)
    process_box(slide, 0.88, 4.02, 3.1, 1.78, "Sparse Retrieval",
                "BM25 captures names, exact terms, programme codes and department keywords.", PALE_GOLD, GOLD)
    process_box(slide, 5.05, 2.42, 2.83, 2.05, "Reciprocal Rank Fusion",
                "Dense and sparse ranked lists are combined into a stronger result list.", WHITE, MAROON)
    process_box(slide, 8.96, 1.84, 3.22, 1.67, "Optional Reranker",
                "Cross-encoder reranker can improve source prioritization.", PALE_BLUE, NAVY)
    process_box(slide, 8.96, 4.12, 3.22, 1.67, "Top-K Evidence",
                "Relevant chunks become context and source cards.", PALE_TEAL, GREEN)
    add_line(slide, 3.98, 2.61, 4.98, 3.25, TEAL, 2)
    add_line(slide, 3.98, 4.84, 4.98, 3.64, GOLD, 2)
    add_line(slide, 7.9, 3.43, 8.9, 2.7, MAROON, 2)
    add_line(slide, 10.57, 3.55, 10.57, 4.05, NAVY, 2)
    add_text(slide, 5.2, 5.2, 2.56, 0.24, "Query expansion supports aliases", 11, MAROON, True, align=PP_ALIGN.CENTER)

    # 11 Answer
    slide = add_slide(prs, 11, "Bilingual and Grounded Answering", "Reliable answers for natural student questions")
    add_rect(slide, 0.82, 1.48, 5.6, 4.95, WHITE, PALE_BLUE, True)
    add_text(slide, 1.08, 1.77, 4.96, 0.26, "English Query", 14, MAROON, True)
    add_rect(slide, 1.08, 2.2, 4.72, 0.68, PALE_BLUE, PALE_BLUE, True)
    add_text(slide, 1.26, 2.42, 4.36, 0.2, "What is the full form of CUSB?", 13, NAVY, True)
    add_rect(slide, 1.55, 3.1, 4.52, 1.02, PALE_TEAL, PALE_TEAL, True)
    add_text(slide, 1.74, 3.35, 4.14, 0.48, "CUSB stands for Central University of South Bihar.", 13, TEAL, True)
    add_text(slide, 1.08, 4.65, 4.98, 0.22, "Source card: About CUSB | official citation metadata", 10, MUTED)
    add_rect(slide, 6.92, 1.48, 5.6, 4.95, WHITE, PALE_BLUE, True)
    add_text(slide, 7.18, 1.77, 4.96, 0.26, "Hinglish Query", 14, MAROON, True)
    add_rect(slide, 7.18, 2.2, 4.72, 0.68, PALE_BLUE, PALE_BLUE, True)
    add_text(slide, 7.36, 2.42, 4.36, 0.2, "CUSB ka full form kya hai?", 13, NAVY, True)
    add_rect(slide, 7.65, 3.1, 4.52, 1.02, PALE_TEAL, PALE_TEAL, True)
    add_text(slide, 7.84, 3.35, 4.14, 0.48, "CUSB ka full form Central University of South Bihar hai.", 13, TEAL, True)
    add_text(slide, 7.18, 4.65, 4.98, 0.22, "Language localization preserves names, values and URLs", 10, MUTED)
    pill(slide, 4.35, 5.62, 4.72, "Direct answer | Grounded LLM | Citation verifier | Fallback", NAVY, WHITE)

    # 12 Stack
    slide = add_slide(prs, 12, "Application Stack", "A deployable research prototype")
    layers = [
        ("Next.js Frontend", "Chat UI, filters, source cards, copy, feedback and mobile layout", MAROON),
        ("FastAPI Backend", "Chat, streaming, search, health, feedback and admin endpoints", TEAL),
        ("RAG Orchestration", "Direct answers, grounded prompt, citation verifier, retries and fallback", NAVY),
        ("Retrieval Layer", "FAISS / local Qdrant option, BM25, query expansion, RRF and reranking", GOLD),
        ("Data and Operations", "Chunk artifacts, embeddings, Docker Compose, analytics and reindexing", GREEN),
    ]
    for i, (head, desc, color) in enumerate(layers):
        y = 1.43 + i * 1.04
        add_rect(slide, 0.92, y, 11.54, 0.75, WHITE, PALE_BLUE, True)
        add_rect(slide, 0.92, y, 0.12, 0.75, color, color)
        add_text(slide, 1.24, y + 0.18, 2.46, 0.24, head, 14, NAVY, True)
        add_text(slide, 3.86, y + 0.18, 8.18, 0.24, desc, 12, MUTED)

    # 13 Evaluation
    slide = add_slide(prs, 13, "Evaluation Design", "Frozen benchmark plus regression testing")
    add_rect(slide, 0.82, 1.42, 3.24, 4.98, PALE_GOLD, PALE_GOLD, True)
    add_text(slide, 1.1, 1.78, 2.7, 0.36, "200 Questions", 27, MAROON, True)
    add_text(slide, 1.1, 2.3, 2.54, 0.38, "Student regression suite", 14, NAVY, True)
    add_text(slide, 1.1, 3.0, 2.52, 1.5,
             "Common student workflows\nPost-fix regression checks\nEnglish and Hinglish pairs\nStable baseline verification",
             13, MUTED)
    add_rect(slide, 4.42, 1.42, 8.1, 4.98, WHITE, PALE_BLUE, True)
    add_text(slide, 4.78, 1.78, 3.9, 0.36, "1,000 Research Questions", 24, TEAL, True)
    cats = [
        "General info", "Admissions", "Departments", "Faculty", "Syllabus",
        "Fees", "Facilities", "Exams / notices", "Bilingual pairs", "Hard negatives",
    ]
    for i, cat in enumerate(cats):
        col, row = i % 2, i // 2
        y = 2.52 + row * 0.65
        x = 4.84 + col * 3.7
        add_circle(slide, x, y, 0.24, TEAL, TEAL)
        add_text(slide, x, y + 0.05, 0.24, 0.14, "100", 7, WHITE, True, align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.36, y + 0.02, 3.06, 0.2, cat, 12, NAVY, True)
    add_text(slide, 4.84, 5.98, 6.98, 0.2, "Robustness: typo | short | vague | privacy | adversarial | future | time-sensitive", 10, MUTED)

    # 14 Metrics
    slide = add_slide(prs, 14, "Evaluation Metrics", "Answer quality, retrieval quality and operational quality")
    blocks = [
        ("Answer Quality", ["Correct", "Partial", "Unsupported", "Hallucinated", "Incomplete", "Error"], MAROON),
        ("Retrieval Quality", ["Recall@5", "MRR", "nDCG@5", "Source relevance", "Citation grounding"], TEAL),
        ("Operational Quality", ["Language match", "Latency", "Timeouts", "Source count", "Retry behavior"], NAVY),
    ]
    for i, (head, items, color) in enumerate(blocks):
        x = 0.82 + i * 4.17
        add_rect(slide, x, 1.57, 3.62, 4.75, WHITE, PALE_BLUE, True)
        add_rect(slide, x, 1.57, 3.62, 0.12, color, color)
        add_text(slide, x + 0.18, 1.95, 3.18, 0.3, head, 16, NAVY, True)
        bullets(slide, items, x + 0.18, 2.55, 3.16, 3.24, 13, MUTED, color, 6)

    # 15 Baseline
    slide = add_slide(prs, 15, "Current 1,000-Question Interim Baseline", "Completed before the latest detector and retry refinements were rerun")
    metric_card(slide, 0.74, 1.41, 2.28, "843", "correct answers", GREEN, "84.3% of benchmark")
    metric_card(slide, 3.2, 1.41, 2.28, "822", "language consistent", TEAL, "82.2% of benchmark")
    metric_card(slide, 5.66, 1.41, 2.28, "492", "grounded citations", MAROON, "49.2% of benchmark")
    metric_card(slide, 8.12, 1.41, 2.08, "0.503", "Recall@5", NAVY, "top-five evidence")
    metric_card(slide, 10.38, 1.41, 2.08, "0.4708", "MRR", GOLD, "first relevant rank")
    bar(slide, 0.84, 3.35, 7.55, "Correct", 843, color=GREEN)
    bar(slide, 0.84, 3.79, 7.55, "Unsupported", 75, color=AMBER)
    bar(slide, 0.84, 4.23, 7.55, "Hallucinated", 66, color=RED)
    bar(slide, 0.84, 4.67, 7.55, "Timeout errors", 12, color=NAVY)
    bar(slide, 0.84, 5.11, 7.55, "Incomplete", 4, color=MUTED)
    add_rect(slide, 10.08, 3.45, 2.32, 1.45, PALE_BLUE, PALE_BLUE, True)
    add_text(slide, 10.3, 3.7, 1.9, 0.28, "nDCG@5", 13, NAVY, True, align=PP_ALIGN.CENTER)
    add_text(slide, 10.3, 4.15, 1.9, 0.35, "0.4763", 27, NAVY, True, align=PP_ALIGN.CENTER)
    add_text(slide, 0.84, 5.93, 11.3, 0.25,
             "Interpretation: a strong prototype baseline with measurable language-routing, hard-negative and source-ranking improvement opportunities.",
             11, MUTED)

    # 16 Latency
    slide = add_slide(prs, 16, "Latency Analysis", "End-to-end API response time from the 1,000-question interim baseline")
    metric_card(slide, 0.84, 1.55, 2.58, "12.97 s", "average latency", MAROON, "affected by timeout outliers")
    metric_card(slide, 3.68, 1.55, 2.58, "4.21 s", "median latency (P50)", TEAL, "typical response time")
    metric_card(slide, 6.52, 1.55, 2.58, "14.65 s", "P95 latency", NAVY, "95% completed below this")
    metric_card(slide, 9.36, 1.55, 2.58, "29.79 ms", "minimum latency", GREEN, "fast direct-answer path")
    add_rect(slide, 0.84, 3.22, 11.1, 1.3, PALE_GOLD, PALE_GOLD, True)
    add_text(slide, 1.12, 3.48, 10.56, 0.75,
             "Latency is measured from API request submission until the complete answer is received. "
             "The evaluator's --delay option is excluded because it only spaces requests to reduce rate-limit errors.",
             15, NAVY, True, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.84, 4.94, 11.1, 0.9, PALE_RED, PALE_RED, True)
    add_text(slide, 1.12, 5.19, 10.56, 0.38,
             "Observed maximum: 4,230.24 s. This extreme timeout/retry outlier explains why median and P95 are reported alongside the mean.",
             12, MAROON, True, align=PP_ALIGN.CENTER)
    add_text(slide, 0.92, 6.28, 10.8, 0.22,
             "Research interpretation: direct answers are fast, while external LLM calls and transient API failures dominate tail latency.",
             11, MUTED)

    # 17 Category
    slide = add_slide(prs, 17, "Category-Level Analysis", "Strengths and improvement areas discovered by testing")
    rows = [
        ("General university info", "100 / 100 stable after refinement", 100, GREEN),
        ("Faculty answers", "100 / 100 answer-quality correct", 100, GREEN),
        ("Syllabus answers", "100 / 100 answer-quality correct", 100, GREEN),
        ("Fee answers", "100 / 100 answer-quality correct", 100, GREEN),
        ("Admission edge cases", "routing refinement needed", 86, AMBER),
        ("Hard-negative handling", "refusal behavior needs strengthening", 24, RED),
    ]
    for i, (label, note, value, color) in enumerate(rows):
        y = 1.52 + i * 0.78
        add_text(slide, 0.85, y, 2.68, 0.22, label, 12, NAVY, True)
        add_rect(slide, 3.66, y + 0.04, 5.0, 0.16, PALE_BLUE, PALE_BLUE, True)
        add_rect(slide, 3.66, y + 0.04, 5.0 * value / 100, 0.16, color, color, True)
        add_text(slide, 8.88, y, 0.72, 0.22, f"{value}%", 12, color, True)
        add_text(slide, 9.62, y, 2.63, 0.22, note, 10, MUTED)
    add_rect(slide, 0.85, 6.3, 11.7, 0.45, PALE_GOLD, PALE_GOLD, True)
    add_text(slide, 1.05, 6.43, 11.25, 0.18,
             "12 reported errors were API timeouts, not confirmed factual failures. Retry support was added after this run.",
             11, MAROON, True)

    # 18 Iteration
    slide = add_slide(prs, 18, "Iterative Improvement Loop", "Evaluation drives scoped engineering changes")
    loop = [
        ("1", "Run", "Frozen benchmark"),
        ("2", "Inspect", "Weak and error rows"),
        ("3", "Cluster", "Failures by intent"),
        ("4", "Fix", "Routing or retrieval"),
        ("5", "Rebuild", "Backend service"),
        ("6", "Rerun", "Affected offsets"),
    ]
    for i, (n, head, desc) in enumerate(loop):
        x = 0.78 + i * 2.06
        add_circle(slide, x + 0.52, 2.0, 0.56, MAROON if i < 3 else TEAL, MAROON if i < 3 else TEAL)
        add_text(slide, x + 0.52, 2.15, 0.56, 0.18, n, 11, WHITE, True, align=PP_ALIGN.CENTER)
        if i < 5:
            add_line(slide, x + 1.12, 2.28, x + 1.96, 2.28, GOLD, 2)
        add_text(slide, x, 2.87, 1.62, 0.24, head, 14, NAVY, True, align=PP_ALIGN.CENTER)
        add_text(slide, x, 3.3, 1.62, 0.36, desc, 11, MUTED, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.1, 4.62, 11.05, 1.02, PALE_TEAL, PALE_TEAL, True)
    add_text(slide, 1.36, 4.84, 10.5, 0.55,
             "Example refinement: English “How do I pay admission fee?” was misclassified because plain “do” was treated as a Hinglish marker. The detector now keeps English “How do ...” queries in English while recognizing Hinglish commands ending in “do”.",
             12, TEAL, True, align=PP_ALIGN.CENTER)

    # 19 Limitations
    slide = add_slide(prs, 19, "Limitations", "A balanced academic assessment")
    limitations = [
        "University information changes over time; official notices remain authoritative.",
        "Direct answers require maintenance and explicit source alignment.",
        "Citation grounding is currently 492/1000 and source ranking needs improvement.",
        "Hard-negative refusal behavior and policy-aware scoring need strengthening.",
        "PDF and OCR extraction may introduce noise.",
        "External LLM quotas and latency affect long batch runs.",
        "A final post-fix full benchmark rerun is still required.",
    ]
    bullets(slide, limitations, 0.94, 1.52, 11.54, 4.95, 16, INK, MAROON, 8)

    # 20 Conclusion
    slide = add_slide(prs, 20, "Conclusion", "What the project demonstrates")
    add_rect(slide, 0.86, 1.43, 11.72, 4.94, NAVY, NAVY, True)
    add_text(slide, 1.24, 1.89, 10.9, 0.82,
             "The project transforms fragmented public CUSB information into a deployable bilingual university assistant.",
             25, WHITE, True, align=PP_ALIGN.CENTER)
    pillars = [
        ("Domain Corpus", "30,384 chunks"),
        ("Hybrid Retrieval", "Dense + BM25 + RRF"),
        ("Grounded Answers", "Citations and fallback"),
        ("Bilingual UX", "English + Hinglish"),
        ("Evaluation", "1,000-question benchmark"),
    ]
    for i, (head, desc) in enumerate(pillars):
        x = 1.2 + i * 2.25
        add_rect(slide, x, 3.44, 1.86, 1.54, WHITE, WHITE, True)
        add_text(slide, x + 0.12, 3.76, 1.62, 0.28, head, 12, NAVY, True, align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.12, 4.34, 1.62, 0.26, desc, 10, MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, 1.18, 5.52, 10.95, 0.26,
             "A strong research prototype with a clear path toward publication-grade evaluation.",
             16, GOLD, True, align=PP_ALIGN.CENTER)

    # 21 Future
    slide = add_slide(prs, 21, "Future Work", "Roadmap toward publication-grade evaluation")
    columns = [
        ("Near-Term Quality", ["Official-source prioritization", "Duplicate-chunk removal", "Source-aligned direct answers", "Stronger refusal handling"], MAROON),
        ("Research Experiments", ["Post-fix 1,000 rerun", "10,000-query stress test", "Baseline comparisons", "Ablation studies"], TEAL),
        ("Deployment Maturity", ["Human annotation", "Hindi support", "Local model fallback", "Scheduled reindexing"], NAVY),
    ]
    for i, (head, items, color) in enumerate(columns):
        x = 0.84 + i * 4.17
        add_rect(slide, x, 1.5, 3.62, 4.95, WHITE, PALE_BLUE, True)
        add_rect(slide, x, 1.5, 3.62, 0.12, color, color)
        add_text(slide, x + 0.18, 1.94, 3.18, 0.28, head, 16, NAVY, True)
        bullets(slide, items, x + 0.18, 2.62, 3.16, 3.1, 13, MUTED, color, 8)

    # 22 References
    slide = add_slide(prs, 22, "References", "Research foundations and official sources")
    references = [
        "Asai, A. et al. (2023). Self-RAG: Learning to Retrieve, Generate, and Critique through Self-Reflection. arXiv:2310.11511.",
        "Es, S. et al. (2023). RAGAS: Automated Evaluation of Retrieval Augmented Generation. arXiv:2309.15217.",
        "Yan, S.-Q. et al. (2024). Corrective Retrieval Augmented Generation. arXiv:2401.15884.",
        "Ru, D. et al. (2024). RAGCHECKER: A Fine-grained Framework for Diagnosing Retrieval-Augmented Generation. arXiv:2408.08067.",
        "Gan, A. et al. (2025). Retrieval Augmented Generation Evaluation in the Era of Large Language Models: A Comprehensive Survey. arXiv:2504.14891.",
        "Guder, L. et al. (2026). RAGe: A Retrieval-Augmented Generation Evaluation Framework. arXiv:2605.27445. Preprint.",
        "Foundations: Lewis et al. (2020) RAG; Karpukhin et al. (2020) DPR; Reimers & Gurevych (2019) Sentence-BERT.",
        "Central University of South Bihar official website: https://www.cusb.ac.in",
    ]
    bullets(slide, references, 0.9, 1.42, 11.65, 5.65, 12, INK, MAROON, 7)

    # 23 Thank you
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_rect(slide, 0, 0, 0.22, 7.5, MAROON, MAROON)
    add_text(slide, 0.92, 1.25, 11.4, 0.5, "THANK YOU", 20, GOLD, True, align=PP_ALIGN.CENTER)
    add_text(slide, 0.92, 2.06, 11.4, 0.95, "Questions and Discussion", 34, WHITE, True, align=PP_ALIGN.CENTER)
    add_text(slide, 1.42, 3.34, 10.4, 0.6,
             "CUSB AI Chatbot\nA Citation-Grounded Bilingual Hybrid RAG System",
             18, RGBColor(220, 230, 239), align=PP_ALIGN.CENTER)
    pill(slide, 4.53, 4.75, 4.27, "Domain Corpus | Hybrid Retrieval | Evaluation", MAROON, WHITE)
    add_text(slide, 11.72, 6.9, 0.8, 0.25, "23", 14, GOLD, True, align=PP_ALIGN.RIGHT)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {len(prs.slides)} slides to {OUT}")


if __name__ == "__main__":
    generate()
